"""
Document Generator Module

Generates PPTX presentations, DOCX documents, and PDF files
from structured outline data. Used by the file_editor skill
for voice-driven document creation.
"""

import os
import subprocess
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH

from core.logger import get_logger


SHARE_DIR = Path(os.path.expanduser("~/jarvis/share"))


class DocumentGenerator:
    """Generates PPTX, DOCX, and PDF documents from structured outlines."""

    def __init__(self, config=None):
        self.logger = get_logger(__name__, config)
        SHARE_DIR.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------------
    # PPTX Generation
    # ------------------------------------------------------------------

    def create_presentation(self, structure: dict, filename: str = "presentation.pptx",
                            images: dict = None) -> Optional[Path]:
        """Generate a PPTX presentation from a structured outline.

        Args:
            structure: Dict with keys: title, subtitle, slides[]
                       Each slide has: title, bullets[], image_query (optional)
            filename: Output filename (saved to share/)
            images: Optional {slide_index: image_path} mapping for embedded images

        Returns:
            Path to saved .pptx file, or None on failure
        """
        try:
            prs = Presentation()
            # Set slide dimensions to widescreen 16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            slides_data = structure.get("slides", [])

            for i, slide_data in enumerate(slides_data):
                slide_title = slide_data.get("title", f"Slide {i + 1}")
                bullets = slide_data.get("bullets", [])
                image_path = images.get(i) if images else None

                if i == 0:
                    # Title slide
                    self._add_title_slide(prs, structure.get("title", slide_title),
                                          structure.get("subtitle", ""))
                elif image_path and Path(image_path).exists():
                    # Content slide with image
                    self._add_image_slide(prs, slide_title, bullets, image_path)
                else:
                    # Content slide (text only)
                    self._add_content_slide(prs, slide_title, bullets)

            # Save
            output_path = SHARE_DIR / filename
            prs.save(str(output_path))
            self.logger.info(f"[doc_gen] Created PPTX: {output_path} ({len(slides_data)} slides)")
            return output_path

        except Exception as e:
            self.logger.error(f"[doc_gen] PPTX creation failed: {e}")
            return None

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Add a title slide (first slide of the deck)."""
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.placeholders[0]
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # Subtitle
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(0x55, 0x55, 0x77)

    def _add_content_slide(self, prs: Presentation, title: str, bullets: list):
        """Add a content slide with title and bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.placeholders[0]
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # Bullets
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for j, bullet in enumerate(bullets):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_after = Pt(8)
                p.level = 0

    def _add_image_slide(self, prs: Presentation, title: str, bullets: list,
                         image_path: str):
        """Add a content slide with text on left and image on right."""
        # Use "Title Only" layout and manually place text + image
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.placeholders[0]
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Text box — left 55%
        text_left = Inches(0.5)
        text_top = Inches(1.8)
        text_width = Emu(int(slide_width * 0.52))
        text_height = Inches(4.5)

        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"\u2022 {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(6)

        # Image — right 40%
        img_left = Emu(int(slide_width * 0.57))
        img_top = Inches(1.8)
        img_width = Emu(int(slide_width * 0.38))
        img_max_height = Inches(4.5)

        try:
            pic = slide.shapes.add_picture(
                str(image_path), img_left, img_top,
                width=img_width
            )
            # Constrain height if needed
            if pic.height > img_max_height:
                ratio = img_max_height / pic.height
                pic.height = img_max_height
                pic.width = int(pic.width * ratio)
        except Exception as e:
            self.logger.warning(f"[doc_gen] Failed to add image to slide: {e}")

    # ------------------------------------------------------------------
    # DOCX Generation
    # ------------------------------------------------------------------

    def create_document(self, structure: dict, filename: str = "document.docx",
                        images: dict = None) -> Optional[Path]:
        """Generate a DOCX document from a structured outline.

        Args:
            structure: Dict with keys: title, subtitle, slides[] (sections)
            filename: Output filename (saved to share/)
            images: Optional {section_index: image_path} mapping

        Returns:
            Path to saved .docx file, or None on failure
        """
        try:
            doc = Document()

            # Document title
            title_para = doc.add_heading(structure.get("title", "Document"), level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Subtitle as italicized paragraph
            subtitle = structure.get("subtitle", "")
            if subtitle:
                sub_para = doc.add_paragraph()
                sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = sub_para.add_run(subtitle)
                run.italic = True
                run.font.size = DocxPt(14)
                run.font.color.rgb = DocxRGB(0x55, 0x55, 0x77)

            doc.add_paragraph()  # Spacer

            sections = structure.get("slides", [])

            for i, section in enumerate(sections):
                # Skip the first "slide" if it's just a title slide with no real content
                if i == 0 and not section.get("bullets"):
                    continue

                section_title = section.get("title", f"Section {i}")
                bullets = section.get("bullets", [])

                # Section heading
                doc.add_heading(section_title, level=1)

                # Image (if available)
                image_path = images.get(i) if images else None
                if image_path and Path(image_path).exists():
                    try:
                        doc.add_picture(str(image_path), width=DocxInches(5.5))
                        last_paragraph = doc.paragraphs[-1]
                        last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    except Exception as e:
                        self.logger.warning(f"[doc_gen] Failed to add image to doc section {i}: {e}")

                # Bullet points
                for bullet in bullets:
                    doc.add_paragraph(bullet, style="List Bullet")

            # Save
            output_path = SHARE_DIR / filename
            doc.save(str(output_path))
            self.logger.info(f"[doc_gen] Created DOCX: {output_path} ({len(sections)} sections)")
            return output_path

        except Exception as e:
            self.logger.error(f"[doc_gen] DOCX creation failed: {e}")
            return None

    # ------------------------------------------------------------------
    # PDF Conversion
    # ------------------------------------------------------------------

    def convert_to_pdf(self, source_path: Path) -> Optional[Path]:
        """Convert a PPTX or DOCX file to PDF via LibreOffice CLI.

        Tries native `libreoffice` command first, then flatpak as fallback.

        Args:
            source_path: Path to the .pptx or .docx file

        Returns:
            Path to the PDF file, or None on failure
        """
        # Try native libreoffice first, then flatpak
        commands = [
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(SHARE_DIR), str(source_path)],
            ["flatpak", "run", "org.libreoffice.LibreOffice",
             "--headless", "--convert-to", "pdf",
             "--outdir", str(SHARE_DIR), str(source_path)],
        ]

        for cmd in commands:
            try:
                result = subprocess.run(
                    cmd, capture_output=True, text=True, timeout=60,
                )
                if result.returncode == 0:
                    pdf_name = source_path.stem + ".pdf"
                    pdf_path = SHARE_DIR / pdf_name
                    if pdf_path.exists():
                        self.logger.info(f"[doc_gen] Converted to PDF: {pdf_path}")
                        return pdf_path
            except FileNotFoundError:
                continue
            except subprocess.TimeoutExpired:
                self.logger.error("[doc_gen] LibreOffice conversion timed out (60s)")
                return None
            except Exception as e:
                self.logger.error(f"[doc_gen] PDF conversion failed: {e}")
                continue

        self.logger.error("[doc_gen] LibreOffice not available — cannot convert to PDF")
        return None
